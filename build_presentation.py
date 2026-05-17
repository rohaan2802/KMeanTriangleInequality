#!/usr/bin/env python3
"""
Build Milestone 3 PowerPoint (.pptx). Requires: python-pptx (use a venv on PEP 668 systems).

Usage (WSL, recommended):
  cd ~/22i-2327_F_FinalProject
  source .venv/bin/activate          # python3 -m venv .venv && pip install python-pptx matplotlib python-docx
  python3 experiments/plot_pdc_figures.py --exp-dir ~/22i-2327_F_FinalProject/experiments
  python3 experiments/build_presentation.py \\
    --figures-dir ~/22i-2327_F_FinalProject/experiments/figures \\
    --out ~/22i-2327_F_FinalProject/PDC_Kmeans_Presentation.pptx

Authors default: read ../AUTHORS.md when --authors omitted. Override: --authors "A (22I-1), B (22I-2)"
"""
from __future__ import annotations

import argparse
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    raise SystemExit("Install: pip install python-pptx (use venv if pip is PEP 668)") from e


def authors_from_readme(figures_dir: Path) -> str:
    """Project root = figures_dir.parent.parent (…/experiments/figures → …/)."""
    root = figures_dir.resolve().parent.parent
    p = root / "AUTHORS.md"
    if not p.is_file():
        return "Mohammad Rohaan (22I-2327) · [Partner name] (22I-xxxx)"
    t = p.read_text(encoding="utf-8", errors="replace")
    m = re.search(r"\|\s*Mohammad Rohaan\s*\|\s*(22I-\d+)\s*\|", t, re.I)
    roll = m.group(1) if m else "22I-2327"
    return f"Mohammad Rohaan ({roll}) · [Partner — edit AUTHORS.md] (22I-xxxx)"


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    sub = slide.placeholders[1]
    sub.text = subtitle


def add_bullets(prs: Presentation, title: str, lines: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)


def add_title_only(prs: Presentation, title: str) -> None:
    layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title


def add_picture_slide(prs: Presentation, title: str, image_path: Path) -> None:
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if image_path.is_file():
        slide.shapes.add_picture(
            str(image_path),
            Inches(0.6),
            Inches(1.35),
            width=Inches(8.9),
        )
    else:
        tx = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        tx.text_frame.text = f"(Add figure: {image_path})"


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument(
        "--figures-dir",
        type=Path,
        default=Path.home() / "22i-2327_F_FinalProject" / "experiments" / "figures",
    )
    ap.add_argument(
        "--out",
        type=Path,
        default=Path.home() / "22i-2327_F_FinalProject" / "PDC_Kmeans_Presentation.pptx",
    )
    ap.add_argument(
        "--authors",
        default=None,
        help='e.g. \'Ali (22I-1234), Sara (22I-5678)\' — if omitted, uses AUTHORS.md next to project root',
    )
    args = ap.parse_args()
    fig = args.figures_dir.expanduser().resolve()
    authors = args.authors if args.authors else authors_from_readme(fig)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Parallel & Distributed Computing — Final Project",
        "Triangle-inequality K-means (author reference code)\n"
        "Gap 1: OpenMP scheduling · Gap 3: Dimensionality\n\n"
        f"{authors}\n"
        "Base paper: Kwedlo & Czochanski, IEEE Access 2019",
    )

    add_bullets(
        prs,
        "Deliverables & artifacts (repo root)",
        [
            "Written report: PDC_Kmeans_Final_Report.docx → export PDF (build_final_report_docx.py).",
            "This deck: PDC_Kmeans_Presentation.pptx (build_presentation.py).",
            "Code + logs: README.md, upstream commit 0265d45a2eb04fec01ed53fc8635b277082ee284.",
            "Figures under experiments/figures/ — regenerate with plot_pdc_figures.py after logs change.",
        ],
    )

    add_bullets(
        prs,
        "Problem & motivation",
        [
            "K-means: alternate assignment (nearest centroid) and update (means).",
            "Lloyd’s assignment is O(N·K·M) per iteration — dominant cost at large N, K.",
            "Triangle-inequality methods (e.g. Elkan) skip many distance evaluations.",
            "Base paper: hybrid MPI/OpenMP implementations on a cluster; scheduling matters.",
        ],
    )

    add_bullets(
        prs,
        "Base paper (short)",
        [
            "Kwedlo & Czochanski — hybrid MPI/OpenMP for Lloyd + Elkan, Annulus, Drake, Yinyang.",
            "Data-parallel assignment; all-reduce for centroid updates.",
            "Compared static vs guided OpenMP on large-scale hardware.",
            "Our work: same codebase, workstation / WSL, two focused gaps.",
        ],
    )

    add_bullets(
        prs,
        "Project scope — two gaps",
        [
            "Gap 1: Static-oriented build vs guided build (DYNAMIC=y) for Elkan.",
            "  Vary OMP_NUM_THREADS on one medium dataset (e.g. N=200k, M=128, K=256).",
            "Gap 3: Vary feature dimension M on synthetic .bin data (N=100k, K=256, T=8).",
            "  Track Avg distance calculations ratio and naive vs Elkan runtime.",
            "Out of scope: multi-node MPI, GPU, modifying authors’ core C++.",
        ],
    )

    add_bullets(
        prs,
        "Methodology",
        [
            "Upstream: hybrid-triangle-kmeans (Bitbucket), fixed commit.",
            "No edits to Clust/*.cpp — two binaries from Makefile flags only:",
            "  • make kmeans OPT=y OPENMP=y  → kmeans_static_gap1",
            "  • make kmeans OPT=y OPENMP=y DYNAMIC=y → kmeans_guided_gap1",
            "Metrics from program output: k-means execution time, iterations, SSE, ratio.",
            "Pairing script: parse_gap1_log.py flags OK vs DIFF_PATH (iter/SSE mismatch).",
        ],
    )

    add_bullets(
        prs,
        "Experimental setup",
        [
            "Environment: WSL Ubuntu, g++, OpenMP (-fopenmp).",
            "K-means: Forgy init, -r 42, -R 1e-4, -a naive or -a elkan.",
            "Gap 1 log: gap1_medium_runs.txt → gap1_parsed_summary.txt",
            "Gap 3 log: gap3_runs.txt (run_gap3_matrix.sh)",
            "Figures: plot_pdc_figures.py → experiments/figures/",
        ],
    )

    add_bullets(
        prs,
        "Live demo (what we can show)",
        [
            "1) Show dataset binary format (optional: ls -lh experiments/data_gap3/).",
            "2) One short run: kmeans_static_gap1 <file.bin> -a elkan -c 256 -r 42 -R 1e-4",
            "3) Point to k-means execution time + Avg distance calculations ratio in stdout.",
            "4) Optional: same command with kmeans_guided_gap1 for contrast.",
        ],
    )

    add_picture_slide(
        prs,
        "Gap 1 — Comparable pairs (median time)",
        fig / "fig_gap1_comparable_pairs_only.png",
    )
    add_picture_slide(
        prs,
        "Gap 1 — All runs median (includes non-comparable T=4,8)",
        fig / "fig_gap1_elkan_static_vs_guided_median.png",
    )

    add_bullets(
        prs,
        "Gap 1 — Interpretation",
        [
            "T = 1,2: pairs usually OK — guided can be slightly faster (load balance).",
            "T = 4,8: often DIFF_PATH — different iterations/SSE between builds.",
            "Cause: parallel float reductions + schedule change trajectory (not a bug).",
            "Fair comparison only when iterations and SSE align.",
        ],
    )

    add_picture_slide(
        prs,
        "Gap 3 — Pruning vs dimension M",
        fig / "fig_gap3_ratio_vs_M.png",
    )
    add_picture_slide(
        prs,
        "Gap 3 — Time vs dimension M",
        fig / "fig_gap3_time_vs_M.png",
    )
    add_picture_slide(
        prs,
        "Gap 3 — Speedup: naive / Elkan (static)",
        fig / "fig_gap3_speedup_naive_over_elkan_static.png",
    )

    add_bullets(
        prs,
        "Gap 3 — Interpretation",
        [
            "Avg distance calculations ratio rises with M → weaker triangle-inequality pruning.",
            "Elkan vs naive speedup shrinks as M grows (uniform synthetic data).",
            "Limitation: i.i.d. uniform is not clustered geometry; trend still illustrates high-D behavior.",
        ],
    )

    add_bullets(
        prs,
        "Limitations & reproducibility",
        [
            "WSL timing noise; some runs show odd thread-clock vs wall-clock lines.",
            "Gap 3: single run per (M, algorithm) in our matrix — medians would be stronger.",
            "Synthetic data only; not the paper’s Tiny Images sample.",
            "README + commit hash + scripts to regenerate logs and figures.",
        ],
    )

    add_bullets(
        prs,
        "Conclusion",
        [
            "Scheduling effects on Elkan are thread- and comparability-dependent.",
            "Dimensionality strongly affects pruning effectiveness (ratio) and Elkan advantage.",
            "Future: repeated trials, clustered synthetic data, deterministic reductions, MPI scale.",
        ],
    )

    add_bullets(
        prs,
        "LLM / AI assistance (course disclosure)",
        [
            "Used for: experiment planning, Makefile/flag explanation, scripts, report/slide structure.",
            "All runs executed locally; results verified (iterations, SSE) where applicable.",
            "Group members understand code paths, metrics, and plots.",
        ],
    )

    add_title_only(prs, "Thank you — Questions?")

    args.out = args.out.expanduser()
    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.out))
    print("Saved:", args.out.resolve())


if __name__ == "__main__":
    main()
